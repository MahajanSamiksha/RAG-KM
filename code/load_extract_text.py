import os
import argparse
import pdfplumber
import pandas as pd
from docx import Document
from pptx import Presentation
from textblob import TextBlob
from gensim.utils import simple_preprocess

# Function to clean text using gensim's simple_preprocess
def clean_text(text):
    return " ".join(simple_preprocess(text, deacc=True))  # Removes punctuation, lowercases

# Function to correct spelling using TextBlob
def correct_spelling(text):
    return str(TextBlob(text).correct())

# Function to extract and preprocess text from a PDF file
def extract_text_from_pdf(pdf_path):
    print(f"extract_text_from_pdf starts: {pdf_path}")
    text = ""
    with pdfplumber.open(pdf_path) as pdf:
        for page in pdf.pages:
            extracted = page.extract_text()
            if extracted:
                text += extracted + "\n"
    print(f"extract_text_from_pdf ends: {pdf_path}")
    #return correct_spelling(clean_text(text.strip()))
    return (clean_text(text.strip()))


# Function to extract and preprocess text from a DOCX file
def extract_text_from_docx(docx_path):
    print(f"extract_text_from_docx starts: {docx_path}")
    doc = Document(docx_path)
    raw_text = "\n".join([para.text for para in doc.paragraphs])    
    print(f"extract_text_from_docx ends: {docx_path}")
    #return correct_spelling(clean_text(raw_text))
    return (clean_text(raw_text))

# Function to extract and preprocess text from an Excel file (XLSX, CSV)
def extract_text_from_excel(file_path):
    print(f"extract_text_from_excel begins: {file_path}")

    text = ""
    if file_path.endswith(".xlsx"):
        df = pd.read_excel(file_path, engine="openpyxl")
    else:  # CSV
        df = pd.read_csv(file_path)
    raw_text = df.to_string(index=False)
    
    print(f"extract_text_from_excel ends: {file_path}")

    #return correct_spelling(clean_text(raw_text))
    return (clean_text(raw_text))

# Function to extract and preprocess text from a PowerPoint file
def extract_text_from_pptx(pptx_path):
    print(f"extract_text_from_pptx begins: {pptx_path}")

    prs = Presentation(pptx_path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    
    print(f"extract_text_from_pptx ends: {pptx_path}")

    #return correct_spelling(clean_text(text.strip()))
    return (clean_text(text.strip()))

# Function to load, extract, clean, and correct text from documents
def load_and_preprocess_text(root_folder):
    extracted_texts = []
    print(f"load_and_preprocess_text begins: {root_folder}")

    for subdir, _, files in os.walk(root_folder):
        print(f"Current Directory: {subdir}")
        print(f"Files: {files}")
        for file in files:
            file_path = os.path.join(subdir, file)
            ext = file.lower().split(".")[-1]
            print(f"load_and_preprocess_text {file_path}")

            try:
                if ext == "pdf":
                    text = extract_text_from_pdf(file_path)
                elif ext == "docx":
                    text = extract_text_from_docx(file_path)
                elif ext in ["xlsx", "csv"]:
                    text = extract_text_from_excel(file_path)
                elif ext == "pptx":
                    text = extract_text_from_pptx(file_path)
                else:
                    continue  # Skip unsupported files
                
                try:
                    if text.strip():
                        print(f"strip is successful")
                        extracted_texts.append({"file": file_path, "text": text})
                        print(f"Processed: {file_path}")
                except Exception as e:
                    print(f"Error processing {file_path}: {str(e)}")

            except Exception as e:
                print(f"Error processing {file_path}: {str(e)}")
    print(f"load_and_preprocess_text ends: {root_folder}")

    return extracted_texts

if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="Count file types and calculate total size in a folder.")
    parser.add_argument("folder", help="Path to the folder (e.g., C:\\Users\\Username\\Documents)")

    args = parser.parse_args()
    folder_path = os.path.abspath(args.folder)

    if not os.path.exists(folder_path):
        print("Error: The specified folder does not exist.")
    else:
        extracted_data = load_and_preprocess_text(folder_path)

    # Save extracted text to a file (optional)
    output_file = "processed_texts.txt"
    with open(output_file, "w", encoding="utf-8") as f:
        for data in extracted_data:
            f.write(f"--- File: {data['file']} ---\n{data['text']}\n\n")

    print(f"\nProcessed text saved to {output_file}")